"""
Generate PowerPoint from HTML Presentation Content
Based on the 20-slide HTML presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    title_shape.text = title
    subtitle_shape.text = subtitle
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)

def add_content_slide(prs, title, content_list, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    title_shape.text = title
    tf = body_shape.text_frame
    tf.clear()
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True

def add_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
    table = slide.shapes.add_table(len(rows)+1, len(headers), x, y, cx, cy).table
    for i in range(len(headers)):
        table.columns[i].width = Inches(9/len(headers))
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 0, 0)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.size = Pt(14)
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True

# SLIDE 1: Title
add_title_slide(prs,
    "Discovering Negative Knowledge in Scientific Literature",
    "An AI-Driven Approach to Identifying Research Gaps Using Graph Neural Networks\n\nDepartment of Computer Science / AI & ML\nAcademic Year: 2025-2026\nFinal Year Capstone Project")

# SLIDE 2: The Problem
add_content_slide(prs, "The Research Challenge", [
    "Information Paradox in Modern Science:",
    "",
    "• 3 million+ research papers published annually",
    "• 50,000+ publications in mental health alone per year",
    "• Researchers cannot comprehensively review all literature",
    "• Valuable interdisciplinary connections remain hidden",
    "• No systematic approach to identify what SHOULD be researched",
    "",
    "Key Question: How do we find what we DON'T know?"
])

# SLIDE 3: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Core Research Question:",
    "How can we systematically discover under-researched connections",
    "in scientific literature using AI?",
    "",
    "Specific Challenges:",
    "• Identifying missing connections between established concepts",
    "• Quantifying confidence in predicted relationships",
    "• Ensuring AI transparency and interpretability",
    "• Scaling across multiple medical domains",
    "• Providing actionable insights for researchers"
])

# SLIDE 4: Research Objectives
add_content_slide(prs, "Research Objectives", [
    "Primary Objective:",
    "Design and develop an AI-powered system that automatically discovers",
    "under-researched connections in scientific literature using Graph Neural Networks",
    "",
    "Key Targets:",
    "• Collect and analyze 500+ research papers",
    "• Extract and classify 600+ mental health concepts",
    "• Construct knowledge graph with 2,000+ relationships",
    "• Achieve >95% prediction accuracy",
    "• Identify 20-50 high-confidence research gaps",
    "• Create interactive 3D visualization platform"
], font_size=16)

# SLIDE 5: Literature Survey Overview
add_content_slide(prs, "Literature Survey: 15 Key Papers Reviewed", [
    "Knowledge Graphs:",
    "  • DBpedia (Auer et al., 2007)",
    "  • YAGO (Rebele et al., 2016)",
    "  • Freebase (Bollacker et al., 2008)",
    "",
    "Biomedical NLP:",
    "  • BioBERT (Lee et al., 2020)",
    "  • ScispaCy (Neumann et al., 2019)",
    "  • PubMedBERT (Gu et al., 2021)",
    "",
    "Graph Neural Networks:",
    "  • Node2Vec (Grover & Leskovec, 2016)",
    "  • GCN (Kipf & Welling, 2017)",
    "  • Link Prediction (Zhang & Chen, 2018)",
    "  • GAT (Veličković et al., 2018)"
], font_size=15)

# SLIDE 6: Research Gaps
add_table_slide(prs, "Identified Research Gaps",
    ["Gap", "Current State", "Impact"],
    [
        ["No automated gap discovery", "Manual expert analysis", "High"],
        ["No domain-specific GNN", "Generic social networks", "High"],
        ["Single data sources", "Limited to one database", "Medium"],
        ["Lack of AI transparency", "Black box systems", "High"],
        ["No confidence metrics", "Binary yes/no suggestions", "High"],
        ["No interactive visualization", "Static reports only", "Medium"],
        ["Limited cross-domain", "Domain-specific only", "Medium"]
    ]
)

# SLIDE 7: System Architecture
add_content_slide(prs, "System Architecture", [
    "End-to-End AI Pipeline:",
    "",
    "DATA COLLECTION → Semantic Scholar + arXiv + PubMed",
    "           ↓",
    "NLP PIPELINE → ScispaCy + BioBERT",
    "           ↓",
    "KNOWLEDGE GRAPH → NetworkX (Concepts + Relations)",
    "           ↓",
    "GRAPH NEURAL NETWORK → GCN for Link Prediction",
    "           ↓",
    "RESEARCH GAP PREDICTION → Confidence Scoring + Ranking"
], font_size=17)

# SLIDE 8: Implementation Phases 1-3
add_content_slide(prs, "Implementation Phases (1-3)", [
    "Phase 1: Data Collection (Weeks 1-2)",
    "  • API integration for 3 databases",
    "  • Collect 500-700 mental health papers",
    "  • Structured database with metadata",
    "",
    "Phase 2: NLP Processing (Weeks 2-3)",
    "  • Entity extraction using ScispaCy and BioBERT",
    "  • Classification: Disorders, Therapies, Risk Factors, Populations, Outcomes",
    "  • Target: 600+ unique concepts",
    "",
    "Phase 3: Relation Extraction (Weeks 3-4)",
    "  • Co-occurrence analysis at sentence level",
    "  • Extract and filter relationships",
    "  • Target: 2,000+ documented relationships"
], font_size=15)

# SLIDE 9: Implementation Phases 4-6
add_content_slide(prs, "Implementation Phases (4-6)", [
    "Phase 4: Knowledge Graph (Week 4)",
    "  • Graph schema design with NetworkX",
    "  • Compute centrality metrics",
    "  • Connected meaningful structure",
    "",
    "Phase 5: Graph Embeddings (Week 5)",
    "  • Node2Vec algorithm implementation",
    "  • 64-dimensional embeddings",
    "  • Similarity validation",
    "",
    "Phase 6: GNN Development (Weeks 6-7)",
    "  • Graph Convolutional Network with PyTorch Geometric",
    "  • Binary cross-entropy loss for link prediction",
    "  • Target: >95% ROC-AUC on validation set"
], font_size=15)

# SLIDE 10: Implementation Phases 7-10
add_content_slide(prs, "Implementation Phases (7-10)", [
    "Phase 7: Gap Prediction (Weeks 7-8)",
    "  • Generate and score missing edges",
    "  • Confidence ranking (0-100%)",
    "  • Select top 20-50 high-confidence gaps",
    "",
    "Phase 8: Visualization (Weeks 8-9)",
    "  • 3D interactive graph with Plotly",
    "  • Transparency panels (architecture, metrics, predictions)",
    "  • User-friendly interface",
    "",
    "Phase 9-10: Evaluation & Documentation (Weeks 9-12)",
    "  • Literature validation and expert review",
    "  • Technical documentation and research paper"
], font_size=15)

# SLIDE 11: Technologies
add_content_slide(prs, "Tools and Technologies", [
    "Programming & Libraries:",
    "  • Python 3.11+",
    "  • SpaCy, scispaCy",
    "  • NetworkX",
    "  • PyTorch + PyTorch Geometric",
    "  • Plotly, Pandas, NumPy",
    "",
    "Data Sources:",
    "  • Semantic Scholar Graph API",
    "  • arXiv API",
    "  • PubMed Entrez E-utilities",
    "",
    "Development:",
    "  • Git version control, Jupyter notebooks, Virtual environments"
], font_size=16)

# SLIDE 12: Expected Outcomes
add_content_slide(prs, "Expected Outcomes", [
    "Technical Deliverables:",
    "  • Functional AI system with all modules integrated",
    "  • Trained GNN model achieving >95% accuracy",
    "  • Interactive 3D visualization platform",
    "  • Knowledge graph of 600+ mental health concepts",
    "",
    "Research Outputs:",
    "  • List of 20-50 high-confidence research gaps",
    "  • Quantitative evaluation metrics",
    "  • Comparison with existing methods",
    "  • System architecture documentation",
    "  • Research methodology paper"
], font_size=16)

# SLIDE 13: Success Criteria
add_table_slide(prs, "Success Criteria",
    ["Criterion", "Target", "Measurement"],
    [
        ["Data Collection", "500+ papers", "Database count"],
        ["Entity Extraction", "600+ concepts", "Entity table count"],
        ["Graph Size", "2,000+ relationships", "Edge count"],
        ["Model Accuracy", ">95% ROC-AUC", "Test set evaluation"],
        ["Literature Validation", "≥70% novel", "Manual verification"],
        ["Expert Rating", "≥4/5 average", "Expert survey"],
        ["Usability", "≥70/100", "SUS questionnaire"]
    ]
)

# SLIDE 14: Project Timeline
add_table_slide(prs, "Project Timeline (12 Weeks)",
    ["Week", "Phase", "Expected Output"],
    [
        ["1-2", "Data Collection", "500+ papers collected"],
        ["2-3", "NLP Processing", "600+ concepts extracted"],
        ["3-4", "Relation Extraction", "2,000+ relationships"],
        ["4-5", "Graph & Embeddings", "Knowledge graph + embeddings"],
        ["6-7", "GNN Training", "Trained model (>95%)"],
        ["7-9", "Prediction & Viz", "20-50 gaps + 3D platform"],
        ["9-12", "Evaluation & Docs", "Final deliverables"]
    ]
)

# SLIDE 15: Resource Requirements
add_content_slide(prs, "Resource Requirements", [
    "Human Resources:",
    "  • Student researcher (full-time, 12 weeks)",
    "  • Faculty advisor (2-3 hours/week)",
    "  • Domain expert (1-2 hours validation)",
    "",
    "Computational Resources:",
    "  • Personal computer (Quad-core, 8GB+ RAM)",
    "  • Internet connection for APIs",
    "  • Google Colab (optional, free GPU)",
    "",
    "Software: All open-source and free",
    "",
    "Estimated Budget: ₹0-5,000 (minimal cost)"
], font_size=17)

# SLIDE 16: Risk Analysis
add_table_slide(prs, "Risk Analysis and Mitigation",
    ["Risk", "Impact", "Mitigation Strategy"],
    [
        ["API rate limits", "Medium", "Exponential backoff; multiple APIs"],
        ["Low NLP accuracy", "High", "Ensemble models; manual validation"],
        ["Sparse graph", "Medium", "Collect more data; augmentation"],
        ["Computational limits", "Medium", "Algorithm optimization; cloud resources"],
        ["Expert validation", "High", "Partner with research groups early"],
        ["Model accuracy", "High", "Try GAT, GraphSAGE alternatives"]
    ]
)

# SLIDE 17: Novel Contributions
add_content_slide(prs, "Novel Contributions", [
    "First Systematic AI Approach to 'Negative Knowledge' Discovery",
    "",
    "Scientific Impact:",
    "  • Transforms serendipitous discovery into systematic methodology",
    "  • Accelerates identification of novel research directions",
    "  • Reduces redundant research efforts",
    "",
    "Technical Innovation:",
    "  • Domain-adapted GNN for medical research",
    "  • Multi-source integration framework",
    "  • AI transparency architecture",
    "",
    "Practical Impact:",
    "  • Actionable gap lists with confidence scores",
    "  • Evidence-based research planning"
], font_size=15)

# SLIDE 18: Future Scope
add_content_slide(prs, "Future Scope", [
    "Expansion Opportunities:",
    "",
    "• Multi-Domain Extension: Cancer, diabetes, cardiovascular research",
    "",
    "• Temporal Analysis: Track research trends and gap evolution over time",
    "",
    "• Causal Inference: Predict causal relationships beyond associations",
    "",
    "• Researcher Networking: Suggest collaboration opportunities",
    "",
    "• Automated Reviews: Generate systematic review sections",
    "",
    "• Database Integration: Direct API with institutional libraries"
], font_size=17)

# SLIDE 19: Key Takeaways
add_content_slide(prs, "Key Takeaways", [
    "The Goal: Transform serendipitous discovery into systematic science",
    "",
    "What Makes This Project Unique?",
    "  • Addresses 7 identified research gaps",
    "  • Combines NLP, Graph Theory, and Deep Learning innovatively",
    "  • Provides complete AI transparency",
    "  • Delivers practical tool with quantified predictions",
    "  • Scalable architecture for multiple medical domains",
    "  • Clear 12-week timeline with achievable milestones",
    "",
    "Impact: Democratize access to research gap analysis and",
    "accelerate scientific discovery"
], font_size=16)

# SLIDE 20: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
left, top, width, height = Inches(1), Inches(2.5), Inches(8), Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Thank You!"
p = tf.paragraphs[0]
p.font.size = Pt(64)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

left2, top2 = Inches(1), Inches(4.5)
txBox2 = slide.shapes.add_textbox(left2, top2, width, Inches(1.5))
tf2 = txBox2.text_frame
tf2.text = "Questions and Discussion\n\nDepartment of Computer Science / AI & ML\nAcademic Year: 2025-2026"
for para in tf2.paragraphs:
    para.font.size = Pt(18)
    para.alignment = PP_ALIGN.CENTER

# Save
prs.save('/Users/srivardhan/Desktop/Research_Proposal_Presentation.pptx')
print("✅ PowerPoint created from HTML presentation!")
print("📁 Saved to: ~/Desktop/Research_Proposal_Presentation.pptx")
print(f"📊 Total slides: {len(prs.slides)}")
