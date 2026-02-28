"""
Generate Comprehensive PowerPoint presentation for project proposal
Professional academic presentation with detailed content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)

def add_content_slide(prs, title, content_list, font_size=16):
    """Add content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]
    
    title_shape.text = title
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content_list:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.space_after = Pt(8)
    
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def add_table_slide(prs, title, headers, rows, col_widths=None):
    """Add slide with table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    
    # Add table
    x, y, cx, cy = Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
    table = slide.shapes.add_table(len(rows)+1, len(headers), x, y, cx, cy).table
    
    # Set column widths
    if col_widths:
        for i, width in enumerate(col_widths):
            table.columns[i].width = Inches(width)
    else:
        for i in range(len(headers)):
            table.columns[i].width = Inches(9/len(headers))
    
    # Add headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.text_frame.paragraphs[0].font.size = Pt(13)
    
    # Add rows
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx+1, col_idx)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    title_shape.text_frame.paragraphs[0].font.size = Pt(30)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

# SLIDE 1: Title
add_title_slide(prs, 
    "Discovering Negative Knowledge in Scientific Literature",
    "An AI-Driven Approach to Research Gap Discovery\nUsing Graph Neural Networks\n\nProject Proposal")

# SLIDE 2: Abstract
add_content_slide(prs, "Abstract", [
    "Scientific publications grow exponentially (3M+ papers/year), creating an information paradox: identifying what we DON'T know becomes harder.",
    "",
    "This project proposes an AI system to discover under-researched connections in mental health literature using Graph Neural Networks and knowledge graphs.",
    "",
    "Approach: Collect 500+ papers from multiple databases → Extract entities via NLP → Build knowledge graph → Train GNN → Predict missing connections with >95% accuracy.",
    "",
    "Impact: Transform manual, serendipitous discovery into systematic, AI-guided research gap identification.",
    "",
    "Deliverables: Functional AI system, 20-50 high-confidence research gaps, interactive 3D visualization, complete documentation."
], font_size=15)

# SLIDE 3: Introduction - The Challenge
add_content_slide(prs, "Introduction: The Research Challenge", [
    "Modern Scientific Landscape:",
    "  • 50,000+ mental health papers published annually",
    "  • Researchers can read ~300 papers/year maximum",
    "  • 99.4% of literature remains unread by any individual",
    "",
    "The Problem:",
    "  • Valuable interdisciplinary connections go undiscovered",
    "  • Duplicate research wastes time and resources",
    "  • Novel research directions remain hidden",
    "",
    "Current State:",
    "  • Tools help find what IS known",
    "  • NO systematic approach to discover what SHOULD BE researched"
], font_size=15)

# SLIDE 4: Introduction - Negative Knowledge Concept
add_content_slide(prs, "Introduction: What is 'Negative Knowledge'?", [
    "Definition:",
    "Negative Knowledge = What we don't know that we don't know",
    "",
    "Examples in Mental Health:",
    "  • Therapy A works for Condition X",
    "  • Therapy A works for Condition Y",
    "  • But has anyone studied Therapy A for Condition Z?",
    "",
    "Research Gap = Missing connection between established concepts",
    "",
    "Challenge: How to systematically identify these gaps?",
    "",
    "Our Solution: Use AI to predict which connections are missing but plausible"
], font_size=16)

# SLIDE 5: Problem Statement
add_content_slide(prs, "Problem Statement", [
    "Core Research Question:",
    "How can we systematically discover under-researched connections in scientific literature using AI?",
    "",
    "Specific Challenges:",
    "1. Information Overload: Cannot manually review all relevant papers",
    "2. Hidden Patterns: Connections span multiple sub-disciplines",
    "3. Lack of Tools: No existing system for gap discovery",
    "4. Confidence Quantification: Need probability scores for predictions",
    "5. Scalability: Must work across multiple medical domains",
    "",
    "Required: Automated, accurate, transparent, and actionable solution"
], font_size=16)

# SLIDE 6: Literature Survey - Introduction
add_content_slide(prs, "Literature Survey: Overview", [
    "Comprehensive review of 15+ research papers across 4 domains:",
    "",
    "1. Knowledge Graph Construction (3 papers)",
    "   → Foundation for representing scientific knowledge",
    "",
    "2. Biomedical NLP (3 papers)",
    "   → Extracting entities from scientific text",
    "",
    "3. Graph Neural Networks (4 papers)",
    "   → Predicting missing connections",
    "",
    "4. Research Gap Identification (5 papers)",
    "   → Current methodologies and limitations",
    "",
    "Goal: Identify gaps our project will address"
], font_size=16)

# SLIDE 7: Literature Survey - Knowledge Graphs
add_table_slide(prs, "Literature Survey: Knowledge Graph Systems",
    ["Paper", "Year", "Key Contribution", "Limitation"],
    [
        ["DBpedia\n(Auer et al.)", "2007", "Automated extraction from Wikipedia\nStructured knowledge base", "General knowledge\nNot domain-specific"],
        ["YAGO\n(Rebele et al.)", "2016", "High-quality temporal KB\nMultilingual support", "Manual curation\nNo predictions"],
        ["Freebase\n(Bollacker et al.)", "2008", "Collaborative graph DB\nLarge-scale coverage", "Limited to explicit\nNo link prediction"]
    ],
    col_widths=[2, 1, 3, 3]
)

# SLIDE 8: Literature Survey - Biomedical NLP
add_table_slide(prs, "Literature Survey: Biomedical NLP Models",
    ["Model", "Year", "Key Achievement", "Gap"],
    [
        ["BioBERT\n(Lee et al.)", "2020", "Pre-trained biomedical BERT\n88%+ NER accuracy", "Entity extraction only\nNo relationships"],
        ["ScispaCy\n(Neumann et al.)", "2019", "Fast biomedical NLP\nMultiple specialized models", "No knowledge graph\nNo gap discovery"],
        ["PubMedBERT\n(Gu et al.)", "2021", "Domain-specific BERT\nPubMed/PMC training", "Lacks prediction\nNo integration"]
    ],
    col_widths=[2, 1, 3.5, 2.5]
)

# SLIDE 9: Literature Survey - Graph Neural Networks
add_table_slide(prs, "Literature Survey: Graph Neural Networks",
    ["Method", "Year", "Innovation", "Application Gap"],
    [
        ["Node2Vec\n(Grover & Leskovec)", "2016", "Scalable embeddings\nRandom walk sampling", "Social networks\nNot scientific lit"],
        ["GCN\n(Kipf & Welling)", "2017", "Convolutional on graphs\nSemi-supervised learning", "Citation graphs\nNo medical domain"],
        ["GAT\n(Veličković et al.)", "2018", "Attention mechanism\nWeighted neighbors", "General graphs\nNo gap discovery"],
        ["Link Prediction\n(Zhang & Chen)", "2018", "SOTA link prediction\n95%+ accuracy", "Social networks\nNot for research"]
    ],
    col_widths=[2, 1, 3, 3]
)

# SLIDE 10: Literature Survey - Research Gap Methods
add_table_slide(prs, "Literature Survey: Gap Identification Methods",
    ["Approach", "Year", "Method", "Critical Limitation"],
    [
        ["Bibliometric\n(Hicks et al.)", "2018", "Citation analysis\nCo-citation patterns", "Manual expert analysis\nNot scalable"],
        ["Keyword-based\n(Wang et al.)", "2019", "Topic modeling\nKeyword gaps", "Misses semantics\nNo confidence scores"],
        ["ABC Model\n(Swanson)", "1997", "Co-occurrence\nHidden connections", "Simple co-occur\nNo deep learning"],
        ["Topic Evolution\n(Various)", "2015-2020", "Temporal analysis\nTrend identification", "Descriptive only\nNo predictions"]
    ],
    col_widths=[2, 1, 3, 3]
)

# SLIDE 11: Research Gaps - Summary Table
add_table_slide(prs, "Identified Research Gaps",
    ["Gap ID", "Research Gap", "Current State", "Our Solution"],
    [
        ["GAP-1", "No Automated\nGap Discovery", "Manual expert\nanalysis required", "Fully automated\nGNN pipeline"],
        ["GAP-2", "No Domain GNN\nfor Science", "GNN for social\nnetworks only", "Medical domain-\nadapted GCN"],
        ["GAP-3", "Single-Source\nData", "One database\n(PubMed/S2)", "Multi-source\naggregation"],
        ["GAP-4", "No AI\nTransparency", "Black box\nsystems", "Complete metric\ndisclosure"],
        ["GAP-5", "No Interactive\nVisualization", "Static reports\nonly", "3D interactive\ngraph"],
        ["GAP-6", "No Confidence\nScores", "Binary yes/no\nsuggestions", "Probabilistic\n0-100% scores"],
        ["GAP-7", "Domain-Specific\nOnly", "Mental health\nonly systems", "Domain-agnostic\narchitecture"]
    ],
    col_widths=[1.2, 2, 2.4, 2.4]
)

# SLIDE 12: Our Contribution
add_content_slide(prs, "Our Unique Contribution", [
    "This project addresses ALL 7 identified gaps:",
    "",
    "Novel Methodology:",
    "  • First AI system specifically for 'negative knowledge' discovery",
    "  • Combines NLP + Knowledge Graphs + GNN innovatively",
    "",
    "Technical Innovation:",
    "  • Domain-adapted GCN for scientific literature (first application)",
    "  • Multi-source data integration (Semantic Scholar + arXiv + PubMed)",
    "  • Probabilistic confidence scoring for research priorities",
    "",
    "Practical Impact:",
    "  • Accelerates gap identification from weeks → hours",
    "  • Provides actionable, ranked research opportunities",
    "  • Complete AI transparency for research validation"
], font_size=15)

# SLIDE 13: Research Objectives
add_content_slide(prs, "Research Objectives", [
    "Primary Objective:",
    "Design and develop an AI system that automatically discovers under-researched connections in scientific literature.",
    "",
    "Specific Objectives:",
    "1. Design multi-source data aggregation pipeline (target: 500+ papers)",
    "2. Develop biomedical NLP pipeline (target: 600+ concepts extracted)",
    "3. Construct knowledge graph (target: 2,000+ documented relationships)",
    "4. Implement & train Graph Convolutional Network (target: >95% ROC-AUC)",
    "5. Generate high-confidence research gap predictions (target: 20-50 gaps)",
    "6. Create interactive 3D visualization with complete AI transparency",
    "7. Validate approach across multiple medical domains"
], font_size=15)

# SLIDE 14: Proposed Methodology - Overview
add_content_slide(prs, "Proposed Methodology: System Overview", [
    "9-Phase AI Pipeline:",
    "",
    "Phase 1-2: Multi-Source Data Collection (500+ papers)",
    "Phase 3: Biomedical NLP Entity Extraction (ScispaCy)",
    "Phase 4: Entity Classification (5 categories)",
    "Phase 5: Relation Extraction (co-occurrence analysis)",
    "Phase 6: Knowledge Graph Construction (NetworkX)",
    "Phase 7: Graph Embedding Training (Node2Vec, 64D)",
    "Phase 8: GNN Model Training (Graph Convolutional Network)",
    "Phase 9: Research Gap Prediction & Visualization",
    "",
    "End-to-end automated pipeline with human validation"
], font_size=15)

# SLIDE 15: Methodology - Phase 1&2: Data Collection
add_content_slide(prs, "Methodology: Data Collection Strategy", [
    "Multi-Source Approach (addresses GAP-3):",
    "",
    "Source 1: Semantic Scholar API",
    "  • 10 search terms × 5 pages × 100 results = 300-500 papers",
    "  • Fields: title, abstract, year, authors, venue",
    "",
    "Source 2: arXiv API",
    "  • Psychology/neuroscience categories (q-bio.NC, cs.CY)",
    "  • 100-150 papers from mental health queries",
    "",
    "Source 3: PubMed Entrez API",
    "  • MeSH term searches with abstract requirement",
    "  • Target: 200+ medical literature papers",
    "",
    "Storage: SQLite database with full metadata tracking"
], font_size=14)

# SLIDE 16: Methodology - Phase 3-5: NLP Pipeline
add_content_slide(prs, "Methodology: NLP Processing Pipeline", [
    "Phase 3: Entity Extraction",
    "  • Models: ScispaCy (en_core_sci_sm + en_ner_bc5cdr_md)",
    "  • Extract: Diseases, chemicals, therapies, populations",
    "  • Confidence: Filter entities with score >0.7",
    "",
    "Phase 4: Entity Classification",
    "  • Rule-based keyword matching into 5 categories:",
    "    - Disorders (depression, anxiety, PTSD...)",
    "    - Therapies (CBT, DBT, medication...)",
    "    - Risk Factors (trauma, stress, abuse...)",
    "    - Populations (adolescent, adult, elderly...)",
    "    - Outcomes (recovery, relapse, mortality...)",
    "",
    "Phase 5: Relation Extraction",
    "  • Co-occurrence at sentence level",
    "  • Create edges between entities appearing together"
], font_size=13)

# SLIDE 17: Methodology - Phase 6-7: Knowledge Graph
add_content_slide(prs, "Methodology: Knowledge Graph Construction", [
    "Phase 6: Graph Building (NetworkX)",
    "  • Nodes: Entities with attributes (name, category, frequency)",
    "  • Edges: Relations with weights (co-occurrence count)",
    "  • Target: 600+ nodes, 2,000+ edges",
    "  • Compute: Degree centrality, betweenness, PageRank",
    "",
    "Phase 7: Graph Embeddings (Node2Vec)",
    "  • Dimensions: 64",
    "  • Walk length: 30 steps",
    "  • Walks per node: 200",
    "  • Window size: 10",
    "  • Output: 64D vector representation for each concept",
    "  • Purpose: Capture semantic similarity for GNN input"
], font_size=14)

# SLIDE 18: Methodology - Phase 8-9: GNN & Prediction
add_content_slide(prs, "Methodology: GNN Training & Gap Prediction", [
    "Phase 8: Graph Convolutional Network",
    "  • Architecture:",
    "    - Input: 64D Node2Vec embeddings",
    "    - Layer 1: GCN(64→64) + ReLU activation",
    "    - Layer 2: GCN(64→32)",
    "    - Decoder: Dot product for link prediction",
    "  • Training: 200 epochs, Adam optimizer (lr=0.01)",
    "  • Loss: Binary cross-entropy",
    "  • Target: >95% ROC-AUC on validation set",
    "",
    "Phase 9: Research Gap Prediction",
    "  • Generate 10,000-15,000 candidate missing edges",
    "  • Score each with trained GNN (0-100% probability)",
    "  • Rank by confidence, select top 20-50",
    "  • Create interactive 3D visualization"
], font_size=13)

# SLIDE 19: Expected Outcomes
add_table_slide(prs, "Expected Outcomes & Deliverables",
    ["Category", "Deliverable", "Specification"],
    [
        ["Technical\nSystem", "Functional AI System", "Multi-source pipeline\nTrained GNN (>95% acc)\n3D visualization"],
        ["Research\nOutputs", "Knowledge Graph", "600+ concepts\n2,000+ connections\nFull metadata"],
        ["Research\nOutputs", "Gap Predictions", "20-50 high-confidence\nRanked by probability\nValidated predictions"],
        ["Documentation", "Complete Docs", "Technical manual\nUser guide\nAPI docs\nPaper draft"],
        ["Academic", "Novel Contribution", "First 'negative knowledge'\nDomain-adapted GNN\nOpen methodology"]
    ],
    col_widths=[1.5, 2.5, 5]
)

# SLIDE 20: Success Metrics
add_table_slide(prs, "Success Criteria & Validation",
    ["Metric", "Target", "Validation Method"],
    [
        ["Data Collection", "500+ papers\n3+ sources", "Database query count\nSource diversity check"],
        ["Entity Extraction", "600+ concepts\n>85% precision", "Manual annotation (50 samples)\nInter-annotator agreement"],
        ["Graph Quality", "2,000+ edges\nConnected graph", "Network analysis\nExpert review of structure"],
        ["Model Accuracy", ">95% ROC-AUC\n>90% Precision@20", "Test set evaluation\nCross-validation"],
        ["Prediction Quality", "20+ gaps >90% conf\n≥70% novel", "Literature verification\nRecent paper check"],
        ["Expert Validation", "≥4/5 rating\n80% plausible", "Domain expert survey\n5-point Likert scale"],
        ["Usability", "≥70/100 SUS score", "User study (5-10 researchers)\nSystem Usability Scale"]
    ],
    col_widths=[2, 2.5, 4.5]
)

# SLIDE 21: Project Timeline
add_table_slide(prs, "12-Week Implementation Timeline",
    ["Week", "Phase", "Key Activities", "Milestone"],
    [
        ["1-2", "Data Collection", "API integration, scraping", "500+ papers"],
        ["2-3", "NLP Processing", "Entity extraction, classification", "600+ concepts"],
        ["3-4", "Relation Extraction", "Co-occurrence, validation", "2,000+ relations"],
        ["4", "Graph Building", "NetworkX, centrality", "Knowledge graph"],
        ["5", "Embeddings", "Node2Vec training", "64D vectors"],
        ["6-7", "GNN Development", "Model design, training", "Trained GNN"],
        ["7-8", "Gap Prediction", "Link prediction, ranking", "20-50 gaps"],
        ["8-9", "Visualization", "3D graph, UI", "Interactive tool"],
        ["9-10", "Evaluation", "Metrics, expert review", "Validation"],
        ["10-12", "Documentation", "Docs, paper writing", "Delivery"]
    ],
    col_widths=[1, 1.8, 3.5, 1.7]
)

# SLIDE 22: Resources Required
add_content_slide(prs, "Resource Requirements", [
    "Human Resources:",
    "  • Student researcher: Full-time (12 weeks)",
    "  • Faculty advisor: 2-3 hours/week for guidance",
    "  • Domain expert: 1-2 hours (validation phase)",
    "",
    "Computational Resources:",
    "  • Personal computer: Quad-core CPU, 8GB+ RAM (sufficient)",
    "  • Google Colab: Optional free GPU for faster training",
    "  • Storage: ~5GB for papers, models, graphs",
    "",
    "Software: All open-source and free",
    "  • Python, PyTorch, NetworkX, scispaCy, Plotly",
    "",
    "Estimated Budget: ₹0-5,000",
    "  • Mostly zero cost (free APIs, open-source tools)",
    "  • Minimal: Expert consultation, cloud computing (if needed)"
], font_size=14)

# SLIDE 23: Risk Mitigation
add_table_slide(prs, "Risk Analysis & Mitigation Strategy",
    ["Risk", "Probability", "Impact", "Mitigation"],
    [
        ["API rate limits", "Medium", "Medium", "Exponential backoff\nMultiple APIs\nLocal caching"],
        ["Low NLP accuracy", "Low", "High", "Ensemble models\nManual validation\nDomain tuning"],
        ["Sparse graph", "Medium", "Medium", "Collect more papers\nGraph augmentation\nLower threshold"],
        ["Model accuracy", "Low", "High", "Try alternatives (GAT)\nHyperparameter tuning\nEnsemble methods"],
        ["Expert access", "Medium", "High", "Early partnerships\nUniversity contacts\nOnline survey option"],
        ["Time constraints", "Low", "Medium", "Phased approach\nCore features first\nExtensions optional"]
    ],
    col_widths=[2, 1.3, 1.3, 4.4]
)

# SLIDE 24: Innovation & Impact
add_content_slide(prs, "Innovation & Societal Impact", [
    "Scientific Innovation:",
    "  • First systematic AI approach to 'negative knowledge' discovery",
    "  • Novel application of GNN to medical literature analysis",
    "  • Methodological contribution to research gap identification",
    "",
    "Practical Impact for Researchers:",
    "  • Accelerate literature review: weeks → hours",
    "  • Discover non-obvious interdisciplinary connections",
    "  • Priority-based research planning with confidence scores",
    "  • Reduce redundant research efforts",
    "",
    "Societal Impact:",
    "  • Faster medical breakthroughs through efficient research",
    "  • Improved mental health treatments and interventions",
    "  • Better allocation of research funding to gaps",
    "  • Democratize access to advanced research tools"
], font_size=14)

# SLIDE 25: Future Scope
add_content_slide(prs, "Future Extensions & Scalability", [
    "Short-term (3-6 months):",
    "  • Expand to diabetes and cancer domains",
    "  • Improve relation typing (causes, treats, associated_with)",
    "  • Add temporal analysis for trend tracking",
    "",
    "Medium-term (6-12 months):",
    "  • Causal inference for mechanism discovery",
    "  • Researcher collaboration network suggestions",
    "  • Integration with institutional research databases",
    "",
    "Long-term Vision:",
    "  • Automated literature review generation",
    "  • Research funding opportunity matching",
    "  • Real-time gap monitoring as papers publish",
    "  • Multi-modal analysis (papers + clinical trials + patents)"
], font_size=15)

# SLIDE 26: Conclusion
add_content_slide(prs, "Conclusion", [
    "This project addresses a fundamental challenge in modern research:",
    "How to systematically discover what we DON'T know",
    "",
    "Key Strengths:",
    "  • Fills 7 identified research gaps in current literature",
    "  • Combines NLP, Knowledge Graphs, and GNN innovatively",
    "  • Clear, achievable 12-week implementation plan",
    "  • Comprehensive validation strategy",
    "  • Potential for significant research impact",
    "",
    "Expected Contribution:",
    "  • Novel methodology for gap discovery",
    "  • Practical tool for researchers",
    "  • Foundation for future research in this area",
    "",
    "Goal: Transform serendipitous discovery into systematic science"
], font_size=15)

# SLIDE 27: Thank You
slide = prs.slides.add_slide(prs.slide_layouts[6])
left = Inches(1)
top = Inches(2.5)
width = Inches(8)
height = Inches(2)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "Thank You!"
p = tf.paragraphs[0]
p.font.size = Pt(64)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 51, 102)
p.alignment = PP_ALIGN.CENTER

# Add contact info
left2 = Inches(1)
top2 = Inches(4.5)
txBox2 = slide.shapes.add_textbox(left2, top2, width, Inches(1.5))
tf2 = txBox2.text_frame
tf2.text = "Questions & Discussion\n\nProject Repository:\ngithub.com/srivardhan-kondu/The-Negative-Knowledge"
for para in tf2.paragraphs:
    para.font.size = Pt(18)
    para.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = '/Users/srivardhan/.gemini/antigravity/brain/ccfe7231-6f67-4621-801e-f6ce52e4beb6/project_proposal_presentation.pptx'
prs.save(output_path)
print("✅ Comprehensive PowerPoint presentation created!")
print(f"📁 Saved to: {output_path}")
print(f"📊 Total slides: {len(prs.slides)}")
print("\n✨ Presentation Features:")
print("  • Abstract slide with full project summary")
print("  • Detailed introduction (2 slides)")
print("  • Comprehensive problem statement")
print("  • In-depth literature survey (5 slides with tables)")
print("  • Clear research gaps analysis (table format)")
print("  • Detailed methodology (6 slides)")
print("  • Complete validation strategy")
print("  • Professional tables and formatting")
